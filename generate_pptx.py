"""Generate Aegis Pitch Deck PPT - Professional Version"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0A, 0x0F, 0x1E)
BG_CARD = RGBColor(0x14, 0x1E, 0x33)
BG_CARD2 = RGBColor(0x1A, 0x25, 0x3D)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0x10, 0xB9, 0x81)
ACCENT3 = RGBColor(0x8B, 0x5C, 0xF6)  # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text="", font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
    return tf


def add_para(tf, text, font_size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(4)
    return p


def add_accent_line(slide, left, top, width, color=ACCENT):
    """Add a colored accent line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_feature_card(slide, left, top, icon_text, title, desc, accent_color=ACCENT):
    """Add a feature card with icon circle"""
    add_rounded_rect(slide, left, top, Inches(3.6), Inches(1.8), BG_CARD, RGBColor(0x2A, 0x3A, 0x55))
    # Icon circle
    circle = add_circle(slide, left + Inches(0.25), top + Inches(0.3), Inches(0.5), accent_color)
    # Icon text in circle
    tf = add_text_box(slide, left + Inches(0.25), top + Inches(0.3), Inches(0.5), Inches(0.5),
                      icon_text, 14, True, WHITE, PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(0.9), top + Inches(0.3), Inches(2.5), Inches(0.4),
                 title, 15, True, WHITE)
    # Description
    add_text_box(slide, left + Inches(0.9), top + Inches(0.75), Inches(2.5), Inches(0.9),
                 desc, 11, False, GRAY)


# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Decorative circles
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x1A, 0x25, 0x50))
add_circle(slide, Inches(10), Inches(5), Inches(4), RGBColor(0x15, 0x20, 0x40))

# Accent line
add_accent_line(slide, Inches(5), Inches(2.8), Inches(3.3), ACCENT)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
             "AEGIS", 80, True, WHITE, PP_ALIGN.CENTER)
tf = add_text_box(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(0.7),
                  "The AI Firewall", 38, False, ACCENT, PP_ALIGN.CENTER)
add_para(tf, "", 8)
add_para(tf, "Enterprise Security & Observability for AI Agents", 20, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# Track badges
add_rounded_rect(slide, Inches(3), Inches(5.2), Inches(3.2), Inches(0.55), BG_CARD, ACCENT)
add_text_box(slide, Inches(3), Inches(5.25), Inches(3.2), Inches(0.5),
             "Track 1: Google AI Studio", 12, True, ACCENT, PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(7), Inches(5.2), Inches(3.5), Inches(0.55), BG_CARD, ACCENT2)
add_text_box(slide, Inches(7), Inches(5.25), Inches(3.5), Inches(0.5),
             "Track 2: Agent Security (Veea)", 12, True, ACCENT2, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.8),
             "Google Gemini 2.0 Flash  +  Veea Lobster Trap  +  Next.js  +  FastAPI", 14, False, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 2: Problem =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), RED)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.5),
             "THE PROBLEM", 13, True, RED)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(1),
             "Enterprises deploy AI agents with", 34, True, WHITE)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(1),
             "ZERO visibility", 34, True, RED)

# Stat cards
stats = [
    ("78%", "of enterprises have AI\nagents in production", ACCENT),
    ("0%", "inspect what goes IN\nor OUT of agents", RED),
    ("$4.5M", "average cost of a\ndata breach (IBM 2025)", YELLOW),
]
for i, (num, desc, color) in enumerate(stats):
    left = Inches(0.8 + i * 4.1)
    add_rounded_rect(slide, left, Inches(3.5), Inches(3.7), Inches(1.8), BG_CARD)
    add_text_box(slide, left + Inches(0.3), Inches(3.7), Inches(3.1), Inches(0.8),
                 num, 36, True, color)
    add_text_box(slide, left + Inches(0.3), Inches(4.5), Inches(3.1), Inches(0.8),
                 desc, 12, False, GRAY)

# Bottom diagram
add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), RGBColor(0x1C, 0x10, 0x10), RED)
add_text_box(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.9),
             "WITHOUT AEGIS:    Your AI Agent  ───────>  LLM API    (completely blind, no inspection)", 16, True, WHITE, PP_ALIGN.CENTER)

# ===== SLIDE 3: Solution =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), ACCENT2)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.5),
             "THE SOLUTION", 13, True, ACCENT2)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1),
             "One-line integration. Total protection.", 34, True, WHITE)

# Flow diagram with cards
# Agent card
add_rounded_rect(slide, Inches(0.5), Inches(3.2), Inches(2.8), Inches(1.2), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.7), Inches(3.4), Inches(2.4), Inches(0.9),
             "Your AI Agent", 16, True, WHITE, PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, Inches(3.4), Inches(3.5), Inches(0.8), Inches(0.5),
             ">>>", 20, True, ACCENT)

# Aegis card (bigger)
add_rounded_rect(slide, Inches(4.2), Inches(2.8), Inches(4.5), Inches(2), BG_CARD, ACCENT2)
add_text_box(slide, Inches(4.4), Inches(2.9), Inches(4.1), Inches(0.5),
             "AEGIS PROXY", 14, True, ACCENT2, PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.4), Inches(3.4), Inches(4.1), Inches(1.3),
             "Deep Prompt Inspection\nReal-time Blocking\nFull Audit Trail", 13, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, Inches(8.8), Inches(3.5), Inches(0.8), Inches(0.5),
             ">>>", 20, True, ACCENT2)

# LLM card
add_rounded_rect(slide, Inches(9.6), Inches(3.2), Inches(3), Inches(1.2), BG_CARD, ACCENT3)
add_text_box(slide, Inches(9.8), Inches(3.4), Inches(2.6), Inches(0.9),
             "Gemini 2.0 Flash", 16, True, WHITE, PP_ALIGN.CENTER)

# 3 Steps
steps = [
    ("1", "Register", "Add your agent (30 sec)"),
    ("2", "Swap URL", "Change LLM base URL to Aegis"),
    ("3", "Monitor", "Real-time events & blocking"),
]
for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.8 + i * 4.1)
    add_rounded_rect(slide, left, Inches(5.5), Inches(3.7), Inches(1.5), BG_CARD)
    add_circle(slide, left + Inches(0.2), Inches(5.7), Inches(0.45), ACCENT)
    add_text_box(slide, left + Inches(0.2), Inches(5.72), Inches(0.45), Inches(0.4),
                 num, 14, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.8), Inches(5.7), Inches(2.7), Inches(0.4),
                 title, 15, True, WHITE)
    add_text_box(slide, left + Inches(0.8), Inches(6.15), Inches(2.7), Inches(0.6),
                 desc, 11, False, GRAY)

# ===== SLIDE 4: Architecture =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), ACCENT3)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.5),
             "ARCHITECTURE & TECH STACK", 13, True, ACCENT3)

# Main architecture box
add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.2), BG_CARD, RGBColor(0x2A, 0x3A, 0x55))
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(6.7), Inches(0.4),
             "Aegis Platform", 16, True, ACCENT2, PP_ALIGN.CENTER)

# Inner boxes
add_rounded_rect(slide, Inches(1.3), Inches(2.6), Inches(6.5), Inches(1.3), RGBColor(0x1A, 0x30, 0x45), ACCENT)
add_text_box(slide, Inches(1.5), Inches(2.7), Inches(6.1), Inches(0.4),
             "Lobster Trap DPI Engine", 14, True, ACCENT)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(6.1), Inches(0.6),
             "13 ingress rules  |  2 egress rules  |  Sub-millisecond  |  Compiled Regex", 11, False, LIGHT_GRAY)

add_rounded_rect(slide, Inches(1.3), Inches(4.1), Inches(3.1), Inches(1), RGBColor(0x1A, 0x30, 0x45), ACCENT2)
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(2.7), Inches(0.35),
             "ALLOW", 13, True, ACCENT2)
add_text_box(slide, Inches(1.5), Inches(4.55), Inches(2.7), Inches(0.4),
             "Forward to Gemini API", 11, False, GRAY)

add_rounded_rect(slide, Inches(4.7), Inches(4.1), Inches(3.1), Inches(1), RGBColor(0x30, 0x15, 0x15), RED)
add_text_box(slide, Inches(4.9), Inches(4.2), Inches(2.7), Inches(0.35),
             "DENY", 13, True, RED)
add_text_box(slide, Inches(4.9), Inches(4.55), Inches(2.7), Inches(0.4),
             "Block + Log + Alert", 11, False, GRAY)

add_rounded_rect(slide, Inches(1.3), Inches(5.3), Inches(6.5), Inches(0.8), RGBColor(0x1A, 0x30, 0x45))
add_text_box(slide, Inches(1.5), Inches(5.4), Inches(6.1), Inches(0.6),
             "FastAPI Backend  |  SQLite  |  SSE Real-time Stream", 12, False, LIGHT_GRAY, PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(1.3), Inches(6.3), Inches(6.5), Inches(0.5), RGBColor(0x1A, 0x30, 0x45))
add_text_box(slide, Inches(1.5), Inches(6.35), Inches(6.1), Inches(0.4),
             "Next.js 16 Dashboard (Real-time SSE)", 12, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# Tech stack on right
stack = [
    ("AI Engine", "Google Gemini 2.0 Flash", ACCENT),
    ("DPI Proxy", "Veea Lobster Trap (Go)", ACCENT2),
    ("Backend", "Python FastAPI + SQLite", ACCENT3),
    ("Frontend", "Next.js 16 + Tailwind", YELLOW),
    ("Deploy", "HF Spaces + Vercel", GRAY),
]
for i, (label, tech, color) in enumerate(stack):
    top = Inches(1.8 + i * 1.05)
    add_rounded_rect(slide, Inches(8.7), top, Inches(4.1), Inches(0.9), BG_CARD)
    add_accent_line(slide, Inches(8.7), top, Pt(4), color)
    add_text_box(slide, Inches(9.0), top + Inches(0.05), Inches(3.6), Inches(0.35),
                 label, 10, True, color)
    add_text_box(slide, Inches(9.0), top + Inches(0.4), Inches(3.6), Inches(0.4),
                 tech, 14, True, WHITE)

# ===== SLIDE 5: Lobster Trap =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), ACCENT2)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.5),
             "VEEA LOBSTER TRAP - DEEP INTEGRATION", 13, True, ACCENT2)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.7),
             "Not just using it - it IS the core engine", 30, True, WHITE)

# Left: Rules
add_rounded_rect(slide, Inches(0.5), Inches(2.6), Inches(6), Inches(4.5), BG_CARD, RGBColor(0x2A, 0x3A, 0x55))
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(5.4), Inches(0.4),
             "13 Ingress Firewall Rules", 15, True, ACCENT2)

rules = [
    "Prompt injection detection",
    "PII / credential leak prevention",
    "SQL injection blocking",
    "Shell command detection",
    "Data exfiltration prevention",
    "Role impersonation blocking",
    "Malware request detection",
    "Obfuscation / encoding detection",
    "Phishing pattern detection",
    "Harmful content blocking",
]
for i, rule in enumerate(rules):
    add_text_box(slide, Inches(1.0), Inches(3.3 + i * 0.37), Inches(5.2), Inches(0.35),
                 f"  {rule}", 11, False, LIGHT_GRAY)
    add_circle(slide, Inches(0.9), Inches(3.38 + i * 0.37), Inches(0.15), ACCENT2)

# Right: Properties
add_rounded_rect(slide, Inches(6.8), Inches(2.6), Inches(6), Inches(4.5), BG_CARD, RGBColor(0x2A, 0x3A, 0x55))
add_text_box(slide, Inches(7.1), Inches(2.8), Inches(5.4), Inches(0.4),
             "Key Properties", 15, True, ACCENT)

props = [
    ("< 1ms", "Inspection latency (compiled regex)"),
    ("0", "LLM calls for security checks"),
    ("15", "Total rules (13 ingress + 2 egress)"),
    ("YAML", "P4-style policy format (open)"),
    ("100%", "Traffic flows through binary"),
]
for i, (val, desc) in enumerate(props):
    top = Inches(3.3 + i * 0.8)
    add_text_box(slide, Inches(7.3), top, Inches(1.5), Inches(0.4),
                 val, 20, True, ACCENT)
    add_text_box(slide, Inches(8.8), top + Inches(0.05), Inches(3.7), Inches(0.4),
                 desc, 12, False, LIGHT_GRAY)

# ===== SLIDE 6: Features =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), ACCENT)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.5),
             "KEY FEATURES", 13, True, ACCENT)

features = [
    ("RT", "Real-time Dashboard", "SSE-powered live security\nevents & metrics streaming", ACCENT),
    ("AG", "Agent Registry", "Register agents, get proxy\nURLs, set policy levels", ACCENT2),
    ("ST", "Security Tester", "Built-in adversarial suite\n27 test cases, 4 categories", ACCENT3),
    ("AU", "Full Audit Trail", "Every request logged with\nmetadata, SOC2-ready", YELLOW),
    ("MT", "Multi-tenant", "User accounts, per-user\nagent isolation & policies", RGBColor(0xEC, 0x48, 0x99)),
    ("ZL", "Zero Latency", "Regex DPI inspection, no\nLLM calls for security", ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.8 + row * 2.8)
    add_feature_card(slide, left, top, icon, title, desc, color)

# ===== SLIDE 7: Business =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), YELLOW)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.5),
             "BUSINESS MODEL & COMPETITIVE EDGE", 13, True, YELLOW)

# Pricing cards
plans = [
    ("Starter", "$99", "/mo", "5 agents\n10K requests\nBasic policies", ACCENT),
    ("Pro", "$299", "/mo", "25 agents\n100K requests\nSlack alerts", ACCENT2),
    ("Enterprise", "Custom", "", "Unlimited agents\nOn-prem deploy\nSOC2, SSO, SLA", ACCENT3),
]
for i, (name, price, period, features, color) in enumerate(plans):
    left = Inches(0.5 + i * 2.8)
    add_rounded_rect(slide, left, Inches(1.7), Inches(2.5), Inches(3.2), BG_CARD, color)
    add_text_box(slide, left + Inches(0.2), Inches(1.85), Inches(2.1), Inches(0.35),
                 name, 12, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.2), Inches(2.1), Inches(0.6),
                 price, 28, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.9), Inches(2.1), Inches(1.5),
                 features, 10, False, GRAY, PP_ALIGN.CENTER)

# Market
add_text_box(slide, Inches(0.5), Inches(5.2), Inches(8), Inches(0.4),
             "TAM: $4.2B (AI security market, projected 2027)", 14, True, YELLOW)
add_text_box(slide, Inches(0.5), Inches(5.7), Inches(8), Inches(0.4),
             "Target: CTOs, Heads of AI, Security Engineers at mid-large enterprises", 12, False, GRAY)

# Competitive edge on right
add_rounded_rect(slide, Inches(8.8), Inches(1.7), Inches(4.2), Inches(5.3), BG_CARD, RGBColor(0x2A, 0x3A, 0x55))
add_text_box(slide, Inches(9.0), Inches(1.85), Inches(3.8), Inches(0.4),
             "Why Aegis Wins", 14, True, WHITE, PP_ALIGN.CENTER)

comparisons = [
    ("Real DPI proxy", "vs webhooks only"),
    ("Sub-ms latency", "vs LLM-based (slow)"),
    ("URL swap (1 line)", "vs SDK required"),
    ("Ingress + Egress", "vs ingress only"),
    ("Open YAML policy", "vs proprietary"),
]
for i, (us, them) in enumerate(comparisons):
    top = Inches(2.4 + i * 0.9)
    add_text_box(slide, Inches(9.2), top, Inches(3.6), Inches(0.35),
                 us, 12, True, ACCENT2)
    add_text_box(slide, Inches(9.2), top + Inches(0.35), Inches(3.6), Inches(0.35),
                 them, 10, False, RED)

# ===== SLIDE 8: Live Demo =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.8), Inches(1.5), ACCENT)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.5),
             "LIVE DEMO", 13, True, ACCENT)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.7),
             "Try it yourself - live right now", 30, True, WHITE)

# URL cards
add_rounded_rect(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2), BG_CARD, ACCENT)
tf = add_text_box(slide, Inches(1.2), Inches(2.75), Inches(2), Inches(0.4),
                  "Frontend", 11, True, GRAY)
add_text_box(slide, Inches(1.2), Inches(3.1), Inches(10.5), Inches(0.5),
             "https://aegis-vert.vercel.app/", 18, True, ACCENT)

add_rounded_rect(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.2), BG_CARD, ACCENT2)
add_text_box(slide, Inches(1.2), Inches(4.15), Inches(2), Inches(0.4),
             "Backend API", 11, True, GRAY)
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(10.5), Inches(0.5),
             "https://rtahreem-aegis-apis.hf.space", 18, True, ACCENT2)

add_rounded_rect(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(0.8), BG_CARD, ACCENT3)
add_text_box(slide, Inches(1.2), Inches(5.5), Inches(5), Inches(0.5),
             "Demo Login:  demo@aegis.ai  /  demo1234", 14, True, ACCENT3)

# Demo steps
steps_text = [
    "1. Login with demo credentials",
    "2. View real-time dashboard (55 sessions, 32 events)",
    "3. Register a new agent - get proxy URL",
    "4. Send adversarial prompts - watch blocks live",
    "5. View security events + audit trail",
]
for i, step in enumerate(steps_text):
    add_text_box(slide, Inches(7), Inches(5.4 + i * 0.38), Inches(6), Inches(0.35),
                 step, 12, False, LIGHT_GRAY)

# ===== SLIDE 9: Thank You =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Decorative
add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(5), RGBColor(0x15, 0x20, 0x40))
add_circle(slide, Inches(10.5), Inches(5), Inches(4), RGBColor(0x15, 0x20, 0x40))

add_accent_line(slide, Inches(5.5), Inches(2.5), Inches(2.3), ACCENT)

add_text_box(slide, Inches(1), Inches(1), Inches(11.3), Inches(1.5),
             "AEGIS", 80, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(0.7),
             "The AI Firewall", 36, False, ACCENT, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
             "One line of code. Total AI agent security.", 20, True, ACCENT2, PP_ALIGN.CENTER)

# Links
add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.4),
             "GitHub: github.com/Rizwantahreem/Aegis", 14, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.9), Inches(11.3), Inches(0.4),
             "Live: https://aegis-vert.vercel.app/", 14, False, GRAY, PP_ALIGN.CENTER)

# Future roadmap mini
add_rounded_rect(slide, Inches(2.5), Inches(5.6), Inches(8.3), Inches(1.4), BG_CARD)
add_text_box(slide, Inches(2.8), Inches(5.7), Inches(7.7), Inches(0.35),
             "ROADMAP", 11, True, ACCENT, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.8), Inches(6.1), Inches(7.7), Inches(0.7),
             "Slack/PagerDuty alerts  |  Visual policy editor  |  SOC2 reports  |  Multi-LLM  |  On-prem", 12, False, GRAY, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11.3), Inches(0.4),
             "Thank you!", 24, True, WHITE, PP_ALIGN.CENTER)

# Save
prs.save("Aegis_Pitch_Deck.pptx")
print("Saved: Aegis_Pitch_Deck.pptx")
